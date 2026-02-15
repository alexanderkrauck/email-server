"""Text extraction from various content types."""

import logging
from typing import Optional, TYPE_CHECKING

if TYPE_CHECKING:
    from src.storage_config.resolver import StorageConfig

logger = logging.getLogger(__name__)


class TextExtractor:
    """Extract text from various content types."""

    async def extract(self, data: bytes, content_type: Optional[str], config: "StorageConfig") -> Optional[str]:
        """
        Extract text from content based on content type.
        
        Args:
            data: Raw binary data
            content_type: MIME type of the content
            config: Storage configuration
            
        Returns:
            Extracted text or None if extraction not supported/configured
        """
        if not data or not content_type:
            return None
            
        from src.storage_config.resolver import should_extract_text
        
        if not should_extract_text(config, content_type):
            return None
            
        content_type_lower = content_type.lower()
        
        try:
            if content_type_lower == "text/plain":
                return self._decode_utf8(data)
            elif content_type_lower in ("text/html", "application/xhtml+xml"):
                return self._extract_html(data)
            elif content_type_lower == "text/csv":
                return self._decode_utf8(data)
            elif content_type_lower == "text/xml":
                return self._decode_utf8(data)
            elif content_type_lower == "application/json":
                return self._decode_utf8(data)
            elif content_type_lower == "application/rtf":
                return self._extract_rtf(data)
            elif content_type_lower == "application/pdf":
                return self._extract_pdf(data)
            elif content_type_lower in (
                "application/msword",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ):
                return self._extract_docx(data)
            elif content_type_lower == "application/vnd.oasis.opendocument.text":
                return self._extract_odt(data)
            elif content_type_lower in (
                "application/vnd.ms-excel",
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ):
                return self._extract_xlsx(data)
            elif content_type_lower == "application/vnd.oasis.opendocument.spreadsheet":
                return self._extract_ods(data)
            elif content_type_lower in (
                "application/vnd.ms-powerpoint",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ):
                return self._extract_pptx(data)
            elif content_type_lower.startswith("image/"):
                return self._extract_image_ocr(data, content_type_lower)
            else:
                logger.debug(f"Unsupported content type for text extraction: {content_type}")
                return None
                
        except Exception as e:
            logger.warning(f"Failed to extract text from {content_type}: {e}")
            return None

    def _decode_utf8(self, data: bytes) -> str:
        """Decode UTF-8 with fallback."""
        return data.decode("utf-8", errors="ignore")

    def _extract_html(self, data: bytes) -> str:
        """Extract text from HTML."""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(data.decode("utf-8", errors="ignore"), "html.parser")
            return soup.get_text(separator=" ", strip=True)
        except Exception as e:
            logger.warning(f"HTML extraction failed: {e}")
            return data.decode("utf-8", errors="ignore")

    def _extract_rtf(self, data: bytes) -> str:
        """Extract text from RTF."""
        try:
            text = data.decode("ascii", errors="ignore")
            text = text.replace("\\par", "\n")
            text = text.replace("\\tab", "\t")
            import re
            text = re.sub(r"\\[a-z]+\d*\s?", "", text)
            text = re.sub(r"[{}]", "", text)
            return text.strip()
        except Exception as e:
            logger.warning(f"RTF extraction failed: {e}")
            return ""

    def _extract_pdf(self, data: bytes) -> str:
        """Extract text from PDF."""
        try:
            import pypdf
            from io import BytesIO
            reader = pypdf.PdfReader(BytesIO(data))
            text_parts = []
            for page in reader.pages:
                text_parts.append(page.extract_text())
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"PDF extraction failed: {e}")
            return ""

    def _extract_docx(self, data: bytes) -> str:
        """Extract text from DOCX."""
        try:
            import docx
            from io import BytesIO
            doc = docx.Document(BytesIO(data))
            text_parts = []
            for para in doc.paragraphs:
                text_parts.append(para.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"DOCX extraction failed: {e}")
            return ""

    def _extract_odt(self, data: bytes) -> str:
        """Extract text from ODT."""
        try:
            import zipfile
            from io import BytesIO
            text_parts = []
            with zipfile.ZipFile(BytesIO(data)) as zf:
                if "content.xml" in zf.namelist():
                    content = zf.read("content.xml").decode("utf-8")
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(content, "xml")
                    for p in soup.find_all("p"):
                        if p.text:
                            text_parts.append(p.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"ODT extraction failed: {e}")
            return ""

    def _extract_xlsx(self, data: bytes) -> str:
        """Extract text from XLSX."""
        try:
            import openpyxl
            from io import BytesIO
            wb = openpyxl.load_workbook(BytesIO(data), data_only=True)
            text_parts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value:
                            text_parts.append(str(cell.value))
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"XLSX extraction failed: {e}")
            return ""

    def _extract_ods(self, data: bytes) -> str:
        """Extract text from ODS."""
        try:
            import zipfile
            from io import BytesIO
            text_parts = []
            with zipfile.ZipFile(BytesIO(data)) as zf:
                if "content.xml" in zf.namelist():
                    content = zf.read("content.xml").decode("utf-8")
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(content, "xml")
                    for cell in soup.find_all("cell"):
                        if cell.text:
                            text_parts.append(cell.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"ODS extraction failed: {e}")
            return ""

    def _extract_pptx(self, data: bytes) -> str:
        """Extract text from PPTX."""
        try:
            from pptx import Presentation
            from io import BytesIO
            prs = Presentation(BytesIO(data))
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            return "\n".join(text_parts)
        except Exception as e:
            logger.warning(f"PPTX extraction failed: {e}")
            return ""

    def _extract_image_ocr(self, data: bytes, content_type: str) -> str:
        """Extract text from images using OCR."""
        try:
            import pytesseract
            from PIL import Image
            from io import BytesIO
            
            image = Image.open(BytesIO(data))
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.warning(f"OCR extraction failed: {e}")
            return ""
